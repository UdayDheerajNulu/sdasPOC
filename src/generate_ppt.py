from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

TITLE = "Database Archival Analysis with LLMs"
SUBTITLE = "Categorization, Archival Columns, and Purge Priorities"

CONTEXT_POINTS = [
	"Business drivers: reduce storage cost, improve performance, ensure compliance",
	"Current state: growing transactional DB, mixed log/config/transaction tables",
	"Constraints: referential integrity, auditability, retention policies",
	"Objective: safe, explainable archival strategy aligned to data lifecycle",
]

PROBLEM_POINTS = [
	"Identify functional groups for existing database tables",
	"Detect archival/retention columns per table",
	"Determine purge order while preserving referential integrity",
	"Automate analysis with consistent, explainable outputs",
]

WORKFLOW_STEPS = [
	"Extract table definitions from SQLite/MySQL",
	"Analyze relationships (FKs, references)",
	"LLM step 1: Categorize tables into functional groups",
	"LLM step 2: Identify primary/secondary archival columns",
	"LLM step 3: Assign intra-group purge priorities",
	"Generate report and Streamlit UI visualization",
]

TECH_STACK = [
	"LangChain (chains, prompts)",
	"LLM: Google Gemini 2.5 Flash (or Groq models)",
	"Databases: SQLite (demo), MySQL",
	"Streamlit (UI)",
	"python-pptx (PPT generation)",
]

def add_title_slide(prs: Presentation):
	slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(silde_layout := slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = TITLE
	subtitle.text = f"{SUBTITLE}\n{datetime.now().strftime('%Y-%m-%d')}"


def add_bulleted_slide(prs: Presentation, heading: str, bullets: list[str]):
	slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(silde_layout := slide_layout)
	title = slide.shapes.title
	title.text = heading
	body = slide.placeholders[1]
	tf = body.text_frame
	tf.clear()
	for i, point in enumerate(bullets):
		p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
		p.text = point
		p.level = 0


def build_presentation(output_path: str = "db_archival_usecase.pptx"):
	prs = Presentation()
	add_title_slide(prs)
	add_bulleted_slide(prs, "Context", CONTEXT_POINTS)
	add_bulleted_slide(prs, "Problem Statement", PROBLEM_POINTS)
	add_bulleted_slide(prs, "Workflow", WORKFLOW_STEPS)
	add_bulleted_slide(prs, "Tech Stack", TECH_STACK)
	prs.save(output_path)
	print(f"✅ PPT generated at: {output_path}")


if __name__ == "__main__":
	build_presentation() 